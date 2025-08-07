import streamlit as st
from pptx import Presentation
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
import os

# Title
st.title("PPTX to PDF Converter (Python pptx)")

# Upload file
uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])

if uploaded_file:
    # Ensure uploads directory exists
    uploads_dir = "uploads"
    os.makedirs(uploads_dir, exist_ok=True)

    # Save uploaded file
    pptx_path = os.path.join(uploads_dir, uploaded_file.name)
    with open(pptx_path, "wb") as f:
        f.write(uploaded_file.read())

    # Load presentation
    presentation = Presentation(pptx_path)

    # Extract text
    all_text = []
    for slide in presentation.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        all_text.append(slide_text.strip())

    # Create PDF
    pdf_filename = uploaded_file.name.replace(".pptx", ".pdf")
    pdf_path = os.path.join(uploads_dir, pdf_filename)

    doc = SimpleDocTemplate(pdf_path)
    styles = getSampleStyleSheet()
    story = []

    for i, text in enumerate(all_text, start=1):
        story.append(Paragraph(f"<b>Slide {i}:</b>", styles["Heading2"]))
        story.append(Paragraph(text.replace("\n", "<br/>"), styles["Normal"]))
        story.append(Spacer(1, 12))

    doc.build(story)

    # Provide download link
    with open(pdf_path, "rb") as f:
        st.download_button("Download PDF", f, file_name=pdf_filename)
